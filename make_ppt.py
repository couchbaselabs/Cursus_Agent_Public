"""Generate Supportal Architecture PPT — Couchbase brand (Open Sans, CB colors)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pathlib import Path

HERE = Path(__file__).parent

# ── Couchbase brand palette ───────────────────────────────────────────────────
C_RED    = RGBColor(0xEC, 0x12, 0x18)
C_DARK   = RGBColor(0x26, 0x26, 0x26)
C_MID    = RGBColor(0x59, 0x59, 0x59)
C_LIGHT  = RGBColor(0xD8, 0xD8, 0xD8)
C_CREAM  = RGBColor(0xFF, 0xF0, 0xDB)
C_BLUE   = RGBColor(0x00, 0xB0, 0xF0)
C_ORANGE = RGBColor(0xFC, 0x9C, 0x0C)
C_GREEN  = RGBColor(0x9D, 0xC2, 0x84)
C_MUTEBL = RGBColor(0x7C, 0xA6, 0xD7)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT  = "Open Sans"
HBAR  = 1.40          # header bar height (inches)
YMAX  = 7.28          # safe bottom edge

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ─────────────────────────────────────────────────────────────────────────────
# Primitives
# ─────────────────────────────────────────────────────────────────────────────
def R(slide, x, y, w, h, fill, border=None, bpt=0.75):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bpt)
    else: s.line.fill.background()
    return s

def T(slide, text, x, y, w, h, sz=14, bold=False, color=C_DARK,
      align=PP_ALIGN.LEFT, italic=False, auto=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    if auto: tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return box

def arrow_h(slide, x1, y, x2):
    cn = slide.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    cn.line.color.rgb = C_LIGHT; cn.line.width = Pt(1.5)

def add_img(slide, path, x, y, w, h):
    p = HERE / path
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        R(slide, x, y, w, h, C_LIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# Composites
# ─────────────────────────────────────────────────────────────────────────────
def header(slide, title, subtitle=""):
    """Red header bar + white body."""
    R(slide, 0, 0, 13.33, HBAR, C_RED)
    R(slide, 0, HBAR, 13.33, 7.5 - HBAR, C_WHITE)
    T(slide, title, 0.30, 0.08, 12.7, 0.78, sz=34, bold=True, color=C_WHITE)
    if subtitle:
        T(slide, subtitle, 0.30, 0.88, 12.7, 0.42,
          sz=13, italic=True, color=RGBColor(0xFF, 0xD8, 0xD0))


TITLE_BAND = 0.46   # height of the colored title band in a card

def card(slide, x, y, w, h, title, bullets, accent=C_BLUE,
         tsz=13, bsz=11):
    """
    Card with full-width accent-colored title band + white body area.
    Text auto-shrinks to fit the body box.
    """
    # outer white box with light border
    R(slide, x, y, w, h, C_WHITE, border=C_LIGHT)
    # title band
    R(slide, x, y, w, TITLE_BAND, accent)
    T(slide, title, x + 0.12, y + 0.06, w - 0.24, TITLE_BAND - 0.08,
      sz=tsz, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # body textbox — auto-shrink to fit
    body_y = y + TITLE_BAND + 0.10
    body_h = h - TITLE_BAND - 0.14
    bx = slide.shapes.add_textbox(
        Inches(x + 0.14), Inches(body_y),
        Inches(w - 0.28), Inches(body_h))
    bx.word_wrap = True
    tf = bx.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2); p.space_after = Pt(0)
        r = p.add_run(); r.text = b
        r.font.name = FONT; r.font.size = Pt(bsz)
        r.font.color.rgb = C_DARK


def flow_box(slide, x, y, w, h, title, body, fill=C_BLUE, tsz=13, bsz=11):
    """Solid-fill box with white text, no sub-structure."""
    R(slide, x, y, w, h, fill)
    T(slide, title, x+0.10, y+0.08, w-0.20, 0.34,
      sz=tsz, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    T(slide, body, x+0.08, y+0.44, w-0.16, h-0.52,
      sz=bsz, color=C_WHITE, align=PP_ALIGN.CENTER, auto=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
R(slide, 0, 0, 13.33, 7.5, C_DARK)
R(slide, 0, 0, 0.40, 7.5, C_RED)
R(slide, 0.40, 5.30, 12.93, 0.08, C_RED)

T(slide, "Supportal AI Platform", 0.75, 1.15, 12.0, 1.15,
  sz=54, bold=True, color=C_WHITE)
T(slide, "Architecture Overview", 0.75, 2.38, 10.0, 0.72,
  sz=30, color=C_RED)
T(slide, "Version 2.6.1  ·  June 2026", 0.75, 3.16, 8.0, 0.44,
  sz=16, color=RGBColor(0xD8, 0xD8, 0xD8))

for i, (name, col, desc) in enumerate([
    ("CURSUS", C_ORANGE, "Launcher & Orchestrator"),
    ("STRABO", C_BLUE,   "NiceGUI Dashboard  :8765"),
    ("CORAX",  C_MUTEBL, "Chainlit AI Chat  :8766"),
]):
    cx = 0.75 + i * 4.15
    R(slide, cx, 3.85, 3.90, 1.00, col)
    T(slide, name, cx+0.18, 3.90, 3.55, 0.46, sz=22, bold=True, color=C_WHITE)
    T(slide, desc, cx+0.18, 4.36, 3.55, 0.36, sz=12, color=C_WHITE)

T(slide, "Couchbase Enterprise 7.6  ·  Python 3.14  ·  Playwright  ·  Claude / OpenAI / Ollama / LMStudio / Gemini",
  0.75, 5.52, 12.0, 0.36, sz=12, color=RGBColor(0x9D, 0x9D, 0x9D))
T(slide, "Confidential — Internal Use Only",
  0.75, 7.06, 6.0, 0.28, sz=10, color=RGBColor(0x59, 0x59, 0x59))


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Names
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "The Names",
       "Each component is named after a figure or concept from the ancient world")

names = [
    ("CURSUS",  C_ORANGE, "docs/ppt_assets/roman_road.jpg",
     "Cursus Publicus",
     "Rome's imperial relay network — a system of way stations and couriers "
     "that kept communications flowing across the entire empire. Cursus launches "
     "Strabo and Corax together and keeps them running, just as the cursus "
     "publicus kept Rome's information moving."),
    ("STRABO",  C_BLUE,   "docs/ppt_assets/strabo_bust.jpg",
     "Strabo of Amaseia  (64 BC – 24 AD)",
     "Greek geographer who wrote Geographica — a 17-volume survey of the known "
     "world. He mapped territories, recorded local customs, and synthesised "
     "everything into a comprehensive reference. Strabo is the dashboard that "
     "maps and visualises your support data."),
    ("CORAX",   C_MUTEBL, "docs/ppt_assets/demosthenes.jpg",
     "Corax of Syracuse  (5th century BC)",
     "Founder of Western rhetoric — the first to write a systematic treatise on "
     "persuasion and argumentation. His name means 'raven' in Greek (κόραξ). "
     "Corax is the AI chat interface that argues from evidence and reasons "
     "over your data."),
]

cw = 4.07
for i, (name, col, img_path, subtitle, body) in enumerate(names):
    cx = 0.28 + i * (cw + 0.14)
    add_img(slide, img_path, cx, HBAR + 0.15, cw, 2.20)
    R(slide, cx, HBAR + 2.35, cw, 0.40, col)
    T(slide, name, cx+0.12, HBAR+2.37, cw-0.24, 0.36,
      sz=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    T(slide, subtitle, cx+0.10, HBAR+2.85, cw-0.20, 0.28,
      sz=11, bold=True, color=col)
    T(slide, body, cx+0.10, HBAR+3.18, cw-0.20, 2.68,
      sz=11, color=C_DARK, auto=True)

R(slide, 0.28, 6.95, 12.77, 0.38, C_DARK)
T(slide, "supportal/  —  The shared Python library connecting all three components. Named for the platform it serves.",
  0.42, 6.98, 12.5, 0.30, sz=12, italic=True, color=C_CREAM,
  align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — System Overview
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "System Overview",
       "Four tiers: Launcher → Applications → Shared Library → Couchbase Storage")

tier_labels = [
    ("TIER 1\nLAUNCHER", C_RED),
    ("TIER 2\nAPPS",      C_RED),
    ("TIER 3\nLIBRARY",  C_RED),
    ("TIER 4\nSTORAGE",  C_RED),
]
tier_ys = [HBAR+0.14, HBAR+1.24, HBAR+2.34, HBAR+3.52]
tier_hs = [0.98, 0.98, 1.06, 1.70]

for (lbl, col), ty, th in zip(tier_labels, tier_ys, tier_hs):
    R(slide, 0.18, ty, 1.18, th, col)
    T(slide, lbl, 0.20, ty+0.14, 1.14, th-0.22,
      sz=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Tier 1 — Cursus
flow_box(slide, 1.50, tier_ys[0], 11.62, 0.98,
         "Cursus  (run_cursus.py)",
         "Single entry-point  ·  Spawns Strabo & Corax as subprocesses  ·  "
         "Color-coded stdout relay  ·  Graceful SIGTERM shutdown  ·  Port overrides via --flags",
         C_ORANGE, tsz=14, bsz=11)

# Tier 2
flow_box(slide, 1.50, tier_ys[1], 5.62, 0.98,
         "Strabo  (apps/strabo/app.py)",
         "NiceGUI  :8765  ·  7 tabs  ·  Scraping, Analytics, Results, Assets, Chat",
         C_BLUE, tsz=12, bsz=11)
flow_box(slide, 7.50, tier_ys[1], 5.62, 0.98,
         "Corax  (apps/corax/app.py)",
         "Chainlit  :8766  ·  Thread persistence  ·  Top-K slider  ·  Shared history",
         C_MUTEBL, tsz=12, bsz=11)
T(slide, "← iframe →", 7.07, tier_ys[1]+0.38, 0.44, 0.22,
  sz=8, color=C_MID, align=PP_ALIGN.CENTER, italic=True)

# Tier 3 — library modules
mods = [
    ("cb_helpers",     "FTS+Vector+N1QL\nRRF · elbow K"),
    ("agent_tools",    "29+ tools\nagent loop"),
    ("prompts",        "System prompt\ndatetime · scope"),
    ("scoring",        "LLM routing\nRAG context"),
    ("prompt_library", "28 prompts\n7 categories"),
    ("llm_providers",  "Claude · OpenAI\nOllama · local"),
]
mw = 11.62 / len(mods)
for i, (m, d) in enumerate(mods):
    col = C_RED if i % 2 == 0 else C_MID
    bx = 1.50 + i * mw
    R(slide, bx, tier_ys[2], mw-0.04, 1.06, col)
    T(slide, m, bx+0.06, tier_ys[2]+0.07, mw-0.14, 0.26,
      sz=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    T(slide, d, bx+0.06, tier_ys[2]+0.36, mw-0.14, 0.66,
      sz=9, color=C_WHITE, align=PP_ALIGN.CENTER, auto=True)

# Tier 4 — Couchbase
R(slide, 1.50, tier_ys[3], 11.62, 1.70, C_CREAM)
R(slide, 1.50, tier_ys[3], 11.62, 0.06, C_RED)
T(slide, "Couchbase Enterprise 7.6  —  localhost:8091  /  couchbase://couchbase (Docker)",
  1.64, tier_ys[3]+0.10, 11.3, 0.30, sz=13, bold=True, color=C_RED)
T(slide, "_default.tickets  ·  _default.snapshots  ·  _default.assets  ·  "
  "chat.history  ·  chat.profiles  ·  chat.threads / steps / elements / users",
  1.64, tier_ys[3]+0.46, 11.3, 0.30, sz=11, color=C_DARK)
T(slide, "8 GSI indexes  ·  FTS index: ticket_fts (BM25 keyword)  ·  "
  "Vector index: ticket_vec (1536-dim cosine)  ·  Idempotent Docker bootstrap",
  1.64, tier_ys[3]+0.80, 11.3, 0.30, sz=11, color=C_DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Cursus detail
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Cursus — The Launcher",
       "run_cursus.py  ·  Named after the cursus publicus, Rome's imperial relay network")

card(slide, 0.28, HBAR+0.18, 6.10, 2.55, "Responsibilities", [
    "• Single entry-point: venv/bin/python run_cursus.py",
    "• Spawns run_strabo.py + run_corax.py as subprocesses",
    "• Color-coded stdout relay: [strabo] blue · [corax] teal",
    "• Watcher daemon threads — if either exits, terminates both",
    "• SIGINT / SIGTERM → clean shutdown of all children",
    "• Port overrides: --strabo-port N  --corax-port N",
], C_ORANGE)

card(slide, 6.68, HBAR+0.18, 6.37, 2.55, "Boot Sequence", [
    "1. run_cursus.py spawns both processes in parallel",
    "2. run_strabo.py → sets STRABO_PORT + PYTHONPATH, execs app.py",
    "3. run_corax.py → sets CORAX_PORT + PYTHONPATH, runs chainlit",
    "4. Strabo → NiceGUI server at http://localhost:8765",
    "5. Corax → Chainlit server at http://localhost:8766",
    "6. Strabo Chat tab renders <iframe src='http://localhost:8766'>",
], C_ORANGE)

# Process diagram
R(slide, 0.28, HBAR+2.90, 12.77, 0.05, C_ORANGE)
T(slide, "Process Model", 0.28, HBAR+2.98, 4, 0.28,
  sz=13, bold=True, color=C_ORANGE)

boxes = [
    (C_ORANGE, "run_cursus.py",       "Process manager\nSIGINT/SIGTERM\nwatcher threads"),
    (C_BLUE,   "run_strabo.py",       "Subprocess 1\nenv: STRABO_PORT\nenv: PYTHONPATH"),
    (C_BLUE,   "apps/strabo/app.py",  "NiceGUI :8765\n~16,700 lines"),
    (C_MUTEBL, "run_corax.py",        "Subprocess 2\nenv: CORAX_PORT\nenv: PYTHONPATH"),
    (C_MUTEBL, "apps/corax/app.py",   "Chainlit :8766\n~900 lines"),
]
bw = 2.18; gap = 0.24
sx = (12.77 - (len(boxes)*bw + (len(boxes)-1)*gap)) / 2 + 0.28
for i, (col, t, b) in enumerate(boxes):
    bx = sx + i*(bw+gap)
    flow_box(slide, bx, HBAR+3.34, bw, 2.05, t, b, col, tsz=12, bsz=10)
    if i < len(boxes)-1:
        labels = ["spawns", "", "execs", ""]
        arrow_h(slide, bx+bw+0.02, HBAR+4.36, bx+bw+gap-0.02)
        if labels[i]:
            T(slide, labels[i], bx+bw, HBAR+4.10, gap+0.04, 0.22,
              sz=8, italic=True, color=C_MID, align=PP_ALIGN.CENTER)

T(slide, "No shared memory between processes — Couchbase is the single shared state store",
  0.28, 7.18, 12.77, 0.24, sz=10, italic=True, color=C_MID, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — Strabo (3-col layout so cards have room)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Strabo — NiceGUI Dashboard",
       "apps/strabo/app.py  ·  Port 8765  ·  ~16,700 lines  ·  7 top-level tabs")

# Use 3 columns, 3 rows to give each card enough height
tabs = [
    ("Configuration", C_MID, [
        "• Auth: cookie paste OR headful SSO (Playwright)",
        "• Couchbase: URL · username · bucket · TLS",
        "• AI Models: Claude · OpenAI · Ollama · LMStudio · Gemini",
        "• Embedding model, batch size, Data Operations",
        "• Preflight: connectivity + index validation",
    ]),
    ("Scraping", C_BLUE, [
        "• Customer directory scrape (Playwright or requests)",
        "• Ticket pipeline: fetch → parse → embed → score → upsert",
        "• Snapshot topology scraping with backfill",
        "• Progress log with cancel support (cancel_event)",
        "• Auto-save to Couchbase after every scrape",
    ]),
    ("Results", C_BLUE, [
        "• Filterable ticket table (priority, status, org)",
        "• Full-text keyword search across subjects",
        "• Ticket detail dialog: description + comments",
        "• Cluster fields: version, RAM, node count",
        "• CSV / Excel export",
    ]),
    ("Chat", C_RED, [
        "• Embeds Corax as full-height <iframe>",
        "• All chat handled by Corax (batch modes retired)",
        "• Consistent agent, tools, and history",
        "• History shared via CB chat.history collection",
        "• CORAX_PORT env var controls iframe src URL",
    ]),
    ("Scoring & Analysis", C_BLUE, [
        "• 12 chart types (bar, line, pie, treemap, scatter…)",
        "• 6 color palettes including Couchbase-branded",
        "• SVG + PNG export per chart; expand/collapse toggle",
        "• Point-click drill-down dialogs (11 chart types)",
        "• Quick Dashboard: 6 pre-built CB metric charts",
    ]),
    ("Customers", C_BLUE, [
        "• Organization directory from Supportal",
        "• Health score badge (≥70 green / ≥40 amber / <40 red)",
        "• Open P1 count column",
        "• Customer switcher → scopes all agent tool calls",
        "• SCOPING RULE: customer= param forced on all queries",
    ]),
    ("Assets", C_ORANGE, [
        "• LLM artifacts stored in CB _default.assets",
        "• Types: report · csv · chart · json · html",
        "• Auto-saved from agent tool calls (background thread)",
        "• Per-asset: Preview · Download · Print · Delete",
        "• Filter by org · type · title keyword",
    ]),
]

cw = 4.22; ch = 2.72; sx = 0.28; sy = HBAR+0.18; gap_x = 0.075
for i, (name, col, bullets) in enumerate(tabs):
    row, c = divmod(i, 3)
    card(slide, sx + c*(cw+gap_x), sy + row*ch, cw, ch-0.10,
         name, bullets, col, tsz=13, bsz=11)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — Corax
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Corax — Chainlit AI Chat",
       "apps/corax/app.py  ·  Port 8766  ·  Thread-persistent  ·  Shared history with Strabo")

card(slide, 0.28, HBAR+0.15, 3.72, 3.00, "Chat Settings  (⚙ sidebar)", [
    "• LLM provider (Claude / OpenAI / Ollama…)",
    "• Model name & API key",
    "• Base URL (local LLMs)",
    "• Top-K retrieval  (slider 1–100, default 10)",
    "• History depth  (slider 2–40 turns)",
    "• Customer scope",
    "• Couchbase connection params",
], C_MUTEBL)

card(slide, 4.22, HBAR+0.15, 4.48, 3.00, "Per-Turn Flow", [
    "1. User sends message",
    "2. Load profile from cl.user_session",
    "3. build_agent_system_prompt() — datetime + scoping",
    "4. call_llm_with_tools() in thread executor",
    "5. Status message updated per tool call (live)",
    "6. Answer parsed for chart/table artifact blocks",
    "7. History → CB chat.history (shared with Strabo)",
], C_MUTEBL)

card(slide, 8.90, HBAR+0.15, 4.15, 3.00, "Session Features", [
    "• Thread persistence (Couchbase-backed)",
    "• Session resume: prior summary injected into prompt",
    "• Follow-up suggestion chips after each answer",
    "• Retry action on tool error",
    "• Scrape job live monitor (asyncio task)",
    "• 28 pre-built prompts (library browser)",
    "• get_current_time → time-aware answers",
], C_MUTEBL)

# Lifecycle hooks
R(slide, 0.28, HBAR+3.30, 12.77, 0.05, C_MUTEBL)
T(slide, "Chainlit Lifecycle Hooks",
  0.28, HBAR+3.40, 5, 0.28, sz=13, bold=True, color=C_MUTEBL)

hooks = [
    ("@cl.on_chat_start",     "Init profile · load history from CB · send welcome (Corax v2.6.1) · configure ChatSettings"),
    ("@cl.on_message",        "Build system prompt → run agent loop in executor → update live status → parse artifacts → save history"),
    ("@cl.on_settings_update","Sync ChatSettings → profile overrides (provider, model, top_k, history_depth, customer…)"),
    ("@cl.action_callback",   "Handles: retry / follow-up chip / rescrape trigger / prompt library selection"),
    ("@cl.on_chat_resume",    "Reload thread from CB · rebuild history · inject prior_session_block into system prompt"),
]
row_h = (YMAX - (HBAR+3.74)) / len(hooks)
for i, (hook, desc) in enumerate(hooks):
    hy = HBAR + 3.74 + i * row_h
    R(slide, 0.28, hy, 3.10, row_h-0.04, C_MUTEBL)
    T(slide, hook, 0.38, hy+0.04, 2.95, row_h-0.08,
      sz=10, bold=True, color=C_WHITE, auto=True)
    T(slide, desc, 3.52, hy+0.04, 9.50, row_h-0.08,
      sz=11, color=C_DARK, auto=True)
    if i < len(hooks)-1:
        R(slide, 0.28, hy+row_h-0.02, 12.77, 0.014, C_LIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — Agent Tool Architecture
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Agent Tool Architecture",
       "call_llm_with_tools()  ·  5 rounds max  ·  OpenAI-compat + native Anthropic  ·  supportal/agent_tools.py")

# Flow steps
steps = [
    (C_ORANGE, "System\nPrompt",   "Today + time\nweekday + week#\nquarter · scope\nTOOL_GUIDANCE"),
    (C_BLUE,   "LLM Call",         "Claude · OpenAI\nOllama · LMStudio\nGemini\nProvider-agnostic"),
    (C_RED,    "Tool Call?",        "Parse tool_calls\n(OpenAI) or\ntool_use\n(Anthropic)"),
    (C_MID,    "_execute_\nagent_tool()", "Dispatch by name\nRuns in process\nReturns string\nresult"),
    (C_GREEN,  "Append\nResult",   "Appended to\nmessage history\nStatus callback\nfired"),
    (C_BLUE,   "Final Answer",     "No more tools\nOR rounds=5\nReturn text\nto caller"),
]
bw = 1.88; gap = 0.28
total_w = len(steps)*bw + (len(steps)-1)*gap
sx = (13.33 - total_w) / 2
box_h = 2.10
for i, (col, t, b) in enumerate(steps):
    bx = sx + i*(bw+gap)
    flow_box(slide, bx, HBAR+0.18, bw, box_h, t, b, col, tsz=12, bsz=10)
    if i < len(steps)-1:
        arrow_h(slide, bx+bw+0.02, HBAR+0.18+box_h/2, bx+bw+gap-0.02)

T(slide, "↺  loop back to LLM Call after each tool result while rounds remaining (max 5)",
  sx, HBAR+2.36, total_w, 0.26, sz=10, italic=True, color=C_MID, align=PP_ALIGN.CENTER)

# Tool catalogue — fits within YMAX
cat_top = HBAR + 2.76
R(slide, 0.28, cat_top, 12.77, 0.05, C_RED)
T(slide, "Tool Catalogue  (29+ tools)",
  0.28, cat_top+0.08, 6, 0.28, sz=13, bold=True, color=C_RED)

cats = [
    ("Query & Search",   C_BLUE,   "query_tickets\ncount_tickets\nhybrid_retrieval\nvector_search_cb"),
    ("Customer Intel",   C_RED,    "get_customer_health_score\ncheck_sla_compliance\nget_digest\nget_portfolio_status"),
    ("Cluster Data",     C_MID,    "get_ticket\nget_cluster_health\nbackfill_snapshot_topology\nget_cluster_risk_report"),
    ("Agent Actions",    C_ORANGE, "tag_ticket\nsave_query\nlist_saved_queries\nsave_artifact"),
    ("Scrape & Refresh", C_GREEN,  "scrape_customer_tickets\nrescrape_customer_tickets\nget_scrape_status\nbackfill_last_comment_at"),
    ("Time & Utility",   C_MUTEBL, "get_current_time ★ NEW\ngenerate_chart\ngenerate_table\ngenerate_customer_report"),
]
cat_content_top = cat_top + 0.42
cat_h = YMAX - cat_content_top
cw = 12.77 / len(cats)
for i, (name, col, tools) in enumerate(cats):
    cx = 0.28 + i*cw
    R(slide, cx, cat_content_top, cw-0.04, cat_h, col)
    s = str(col); dr = RGBColor(
        max(0, int(s[0:2],16)-30), max(0, int(s[2:4],16)-30), max(0, int(s[4:6],16)-30))
    R(slide, cx, cat_content_top, cw-0.04, 0.32, dr)
    T(slide, name, cx+0.08, cat_content_top+0.04, cw-0.16, 0.26,
      sz=10, bold=True, color=C_WHITE)
    T(slide, tools, cx+0.08, cat_content_top+0.36, cw-0.16, cat_h-0.42,
      sz=10, color=C_WHITE, auto=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — supportal Library
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "supportal/  — Shared Intelligence Library",
       "Extracted from monolith in v2.2–v2.4  ·  −5,600 lines from strabo/app.py  ·  imported by both Strabo and Corax")

modules = [
    ("cb_helpers.py",      C_RED,    [
        "• Couchbase SDK 4.x connection management",
        "• load_to_couchbase() + cancel_event support",
        "• hybrid_retrieval() — FTS + Vector + N1QL → RRF",
        "• _rrf_elbow_k() — adaptive Top-K via elbow detection",
        "• Asset / history / profile CRUD helpers",
    ]),
    ("agent_tools.py",     C_BLUE,   [
        "• _AGENT_TOOLS (29+ OpenAI-format tool defs)",
        "• call_llm_with_tools() — 5-round agent loop",
        "• Fleet analytics: portfolio, health, SLA, digest",
        "• Asset management: save / list / get / delete",
        "• get_current_time (new v2.6.1)",
    ]),
    ("prompts.py",         C_ORANGE, [
        "• build_agent_system_prompt()",
        "• Injects: date + time + weekday + week + quarter",
        "• Customer scoping rule injection",
        "• TOOL_GUIDANCE per-tool instructions",
        "• Prior session block injection",
    ]),
    ("scoring.py",         C_MID,    [
        "• LLM provider routing",
        "• RAG context builder",
        "• score_tickets_batch()",
        "• Complexity (1–10) + sentiment scoring",
        "• Cluster ↔ app alias maps",
    ]),
    ("prompt_library.py",  C_GREEN,  [
        "• 28 pre-built analyst prompts",
        "• 7 categories (health / SLA / risk…)",
        "• {customer} injection at runtime",
        "• Chainlit two-step browser widget",
        "• NiceGUI nested expansion panels",
    ]),
    ("llm_providers.py",   C_MUTEBL, [
        "• call_llm() — unified single-call interface",
        "• Claude (native Anthropic SDK tool use)",
        "• OpenAI-compat (OpenAI SDK)",
        "• Ollama / LMStudio (local base_url)",
        "• Gemini (openai-compat endpoint)",
    ]),
    ("couchbase_data_layer.py", C_RED, [
        "• CouchbaseDataLayer class",
        "• High-level ORM-like wrapper",
        "• Used by Corax for settings persistence",
        "• Ticket / snapshot / profile CRUD",
    ]),
    ("snapshot_parser.py\nticket_parser.py", C_BLUE, [
        "• HTML → structured dict parsers",
        "• Bucket: mem, eviction, storage_mode",
        "• GSI / FTS / Eventing counts + names",
        "• N2N encryption, RAM quotas",
        "• Scope / collection counts per bucket",
    ]),
]

cols = 4; cw = 3.20; ch = 2.72; sx = 0.28; gap_x = 0.075
for i, (name, col, bullets) in enumerate(modules):
    row, c = divmod(i, cols)
    card(slide, sx + c*(cw+gap_x), HBAR+0.18 + row*ch,
         cw, ch-0.10, name, bullets, col, tsz=12, bsz=11)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — Data Pipeline
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Data Pipeline",
       "Scrape → Parse → Embed → Score → Store → Query  ·  Idempotent upserts  ·  cancel_event support")

# Stage flow
stages = [
    (C_ORANGE, "Scrape",   "Playwright / requests\nSupportal HTML\nTickets + Snapshots\nZendesk comments"),
    (C_BLUE,   "Parse",    "ticket_parser.py\nsnapshot_parser.py\nStructured dicts\nCluster topology"),
    (C_MUTEBL, "Embed",    "OpenAI text-embed-3\n1536-dim vectors\nStored on ticket\ndocument"),
    (C_RED,    "Score",    "scoring.py\nComplexity 1–10\nSentiment\nPriority risk"),
    (C_MID,    "Store",    "load_to_couchbase()\nBatch upsert\ncancel_event\nIdempotent keys"),
    (C_GREEN,  "Query",    "hybrid_retrieval()\nFTS + Vector\n+ N1QL → RRF\nElbow Top-K"),
]
bw = 1.88; gap = 0.28
total_w = len(stages)*bw + (len(stages)-1)*gap
sx = (13.33-total_w)/2
stage_h = 2.20
for i, (col, t, b) in enumerate(stages):
    bx = sx + i*(bw+gap)
    flow_box(slide, bx, HBAR+0.18, bw, stage_h, t, b, col, tsz=13, bsz=10)
    if i < len(stages)-1:
        arrow_h(slide, bx+bw+0.02, HBAR+0.18+stage_h/2, bx+bw+gap-0.02)

# RRF detail section
rrf_top = HBAR + 2.54
R(slide, 0.28, rrf_top, 12.77, 0.05, C_RED)
T(slide, "Hybrid Retrieval  —  cb_helpers.hybrid_retrieval()",
  0.28, rrf_top+0.08, 8, 0.28, sz=13, bold=True, color=C_RED)

rrf_cols = [
    ("FTS Search\n(BM25 keyword)",    C_BLUE,
     "fts_keyword_search_cb()\nstruct_keywords field\nsubject + description\ncomments text"),
    ("Vector Search\n(cosine sim)",   C_MUTEBL,
     "vector_search_cb()\nembedding field\n1536-dim query embed\nANN nearest-neighbor"),
    ("N1QL Structured\n(exact match)",C_RED,
     "Priority / status\nDate range predicates\nOrg scoping\nSort by last_comment"),
    ("RRF Fusion\n(Reciprocal Rank)", C_ORANGE,
     "score = Σ 1/(k + rank)\nMerges all three lists\nDe-dup by ticket_id\nreciprocal_rank_fusion()"),
    ("Elbow Top-K\n(adaptive cutoff)",C_GREEN,
     "_rrf_elbow_k()\nLargest relative drop\nbetween consecutive\nRRF scores"),
]
rrf_content_top = rrf_top + 0.42
rrf_h = YMAX - rrf_content_top
rcw = 12.77 / len(rrf_cols)
for i, (name, col, body) in enumerate(rrf_cols):
    cx = 0.28 + i*rcw
    R(slide, cx+0.04, rrf_content_top, rcw-0.10, rrf_h, C_WHITE, border=col)
    R(slide, cx+0.04, rrf_content_top, rcw-0.10, 0.05, col)
    T(slide, name, cx+0.12, rrf_content_top+0.08, rcw-0.24, 0.38,
      sz=11, bold=True, color=col)
    T(slide, body, cx+0.12, rrf_content_top+0.50, rcw-0.24, rrf_h-0.58,
      sz=11, color=C_DARK, auto=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 10 — Couchbase Backend
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Couchbase Backend",
       "Enterprise 7.6.6 build 6126  ·  Bucket: supportal  ·  8 GSI + FTS + Vector indexes")

# Left — collections (compact)
T(slide, "Collections", 0.28, HBAR+0.18, 3.5, 0.28,
  sz=14, bold=True, color=C_RED)

scopes = [
    ("_default scope", C_BLUE, [
        ("tickets",   "ticket::<id>",           "Main docs — embedding, scores, comments, struct_keywords"),
        ("snapshots", "snapshot::<uuid>::<ver>", "Cluster topology — buckets, indexes, nodes, CBSE"),
        ("assets",    "asset::<session>::<uuid>","LLM artifacts — reports, CSV, charts, JSON, HTML"),
    ]),
    ("chat scope", C_MUTEBL, [
        ("history",       "history::<org>::<ts>", "Per-customer chat log shared Strabo ↔ Corax"),
        ("profiles",      "profile::<username>",  "Config: provider, model, top_k, customer scope"),
        ("threads/steps", "Chainlit thread meta", "Session metadata, message history, sidebar"),
    ]),
]

row_h_coll = 0.62
cy = HBAR + 0.52
for scope_name, col, colls in scopes:
    R(slide, 0.28, cy, 6.22, 0.30, col)
    T(slide, scope_name, 0.42, cy+0.04, 5.9, 0.24,
      sz=11, bold=True, color=C_WHITE)
    cy += 0.30
    for cname, key, desc in colls:
        R(slide, 0.28, cy, 6.22, row_h_coll, C_WHITE, border=C_LIGHT)
        R(slide, 0.28, cy, 0.046, row_h_coll, col)
        T(slide, cname, 0.40, cy+0.06, 1.55, 0.26, sz=11, bold=True, color=C_DARK)
        T(slide, key,   0.40, cy+0.32, 5.7, 0.20, sz=8, italic=True, color=C_MID)
        T(slide, desc,  2.08, cy+0.06, 4.3, 0.52, sz=10, color=C_DARK, auto=True)
        cy += row_h_coll
    cy += 0.10

# Right — indexes
T(slide, "Indexes", 6.78, HBAR+0.18, 6.5, 0.28,
  sz=14, bold=True, color=C_RED)

indexes = [
    (C_BLUE,   "GSI Indexes  (N1QL structured queries) — 8 total",
     "idx_tickets_org · idx_tickets_status · idx_tickets_priority\n"
     "idx_tickets_created · idx_tickets_updated · idx_last_comment\n"
     "idx_tickets_org_status · idx_assets_org\n"
     "+ Primary indexes on _default.assets and chat.profiles"),
    (C_RED,    "FTS Index  —  ticket_fts  (keyword / BM25)",
     "Fields: subject, description, comments[], struct_keywords\n"
     "Analyzer: standard  ·  Store: true (snippet highlighting)\n"
     "Used by: fts_keyword_search_cb() in hybrid_retrieval()"),
    (C_GREEN,  "Vector Index  —  ticket_vec  (semantic similarity)",
     "Field: embedding  ·  Dimensions: 1536  ·  Similarity: cosine\n"
     "Model: text-embedding-3-small (OpenAI) or local equivalent\n"
     "Used by: vector_search_cb() in hybrid_retrieval()"),
    (C_ORANGE, "RRF Fusion  (Python — no additional CB index needed)",
     "Merges FTS + vector + N1QL ranked lists\n"
     "_rrf_elbow_k(): natural cutoff via largest relative score drop\n"
     "Top-K slider = hard upper cap; elbow can only reduce"),
]
idx_top = HBAR + 0.52
idx_h = (YMAX - idx_top - 0.16) / len(indexes)
for i, (col, title, desc) in enumerate(indexes):
    iy = idx_top + i * idx_h
    R(slide, 6.78, iy, 6.27, idx_h-0.06, C_WHITE, border=C_LIGHT)
    R(slide, 6.78, iy, 0.046, idx_h-0.06, col)
    T(slide, title, 6.92, iy+0.08, 6.0, 0.26, sz=11, bold=True, color=col)
    T(slide, desc,  6.92, iy+0.38, 6.0, idx_h-0.50, sz=10, color=C_DARK, auto=True)

T(slide, "⚠  scoring_model is not a valid FTS index key in CB 7.6 — omit from all FTS index definitions",
  6.78, YMAX+0.02, 6.27, 0.22, sz=8, italic=True, color=C_MID)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 11 — Docker & Deployment
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Docker & Deployment",
       "docker compose up --build  ·  Idempotent Couchbase bootstrap  ·  Single-command start")

card(slide, 0.28, HBAR+0.18, 6.18, 2.62, "docker-compose.yml Services", [
    "couchbase      — CB Enterprise/Community image, ports 8091-8096 + 11210",
    "couchbase-init — one-shot (restart: 'no'), mounts couchbase-init.sh",
    "               depends_on: couchbase (health check passes first)",
    "app            — Cursus launcher",
    "               depends_on: couchbase-init: service_completed_successfully",
    "docker compose down -v  →  full data reset",
], C_BLUE)

card(slide, 6.74, HBAR+0.18, 6.31, 2.62, "couchbase-init.sh  (pure curl, idempotent)", [
    "1. Poll /ui/index.html until CB REST API responds",
    "2. Check /pools/default — skip if already initialized",
    "3. cluster-init: RAM quotas, services (kv/n1ql/index/fts)",
    "4. bucket-create: supportal + CB_BUCKET_RAMSIZE",
    "5. Create scopes (_default exists) and chat",
    "6. Create collections: tickets · snapshots · assets · history · profiles",
    "7. POST 8 GSI + primary index statements via REST API",
], C_BLUE)

card(slide, 0.28, HBAR+3.00, 4.08, 2.38, "Local Development", [
    "python3 -m venv venv",
    "source venv/bin/activate",
    "pip install -r requirements.txt",
    "playwright install chromium",
    "",
    "venv/bin/python run_cursus.py",
    "  → Strabo:  http://localhost:8765",
    "  → Corax:   http://localhost:8766",
], C_MID)

card(slide, 4.56, HBAR+3.00, 4.10, 2.38, "Key Environment Variables", [
    "SUPPORTAL_URL        Supportal base URL",
    "CB_URL / CB_USER / CB_PASS / CB_BUCKET",
    "CB_RAMSIZE / CB_INDEX_RAMSIZE / CB_FTS_RAMSIZE",
    "ANTHROPIC_API_KEY  /  OPENAI_API_KEY",
    "STRABO_PORT  (default 8765)",
    "CORAX_PORT   (default 8766)",
    "CHAINLIT_AUTH_SECRET",
], C_MID)

card(slide, 8.96, HBAR+3.00, 4.09, 2.38, "Access Points", [
    "Strabo:    http://localhost:8765",
    "Corax:     http://localhost:8766",
    "CB Admin:  http://localhost:8091",
    "",
    "Config → Couchbase:",
    "  URL:  couchbase://couchbase (Docker)",
    "  Bucket: supportal → Save & Test",
], C_ORANGE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 12 — Roadmap Phase 3
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Roadmap  —  Phase 3: Fleet Analysis",
       "Shift from single-customer interrogation to fleet-wide intelligence  ·  docs/phase3-fleet-analysis.md")

milestones = [
    ("3.1  Fleet Query Foundation", C_RED, [
        "query_fleet_tickets — cross-org N1QL with GROUP BY",
        "list_at_risk_clusters — bad/warn + no open ticket",
        "fleet_version_distribution — CB version counts fleet-wide",
        "fleet_cbse_impact — CBSEs ranked by org blast radius",
        "get_portfolio_status + cluster bad_ratio dimension",
    ]),
    ("3.2  Fleet Dashboard Tab", C_BLUE, [
        "New Fleet tab (between Customers and Assets)",
        "Charts: version donut · tickets by org h-bar",
        "Priority stacked bar · cluster bad-item heatmap",
        "30-day ticket trend area chart",
        "KPI chips: total open · orgs with P1 · bad clusters",
    ]),
    ("3.3  Leading Indicators", C_ORANGE, [
        "detect_leading_indicators — risk score formula",
        "risk = (bad×3 + warn) × recency_factor",
        "fleet_anomaly_scan — >2σ vs 7-day rolling baseline",
        "Alerts panel in Fleet Dashboard",
        "get_cluster_risk_report — full health history + tickets",
    ]),
    ("3.4  Portfolio Management", C_GREEN, [
        "saved_portfolio::{name} docs in CB tickets collection",
        "create_portfolio · list_portfolios · get_portfolio_health",
        "Aggregate score + SLA compliance + open P1 count",
        "Fleet Dashboard portfolio switcher filters all charts",
        "Portfolio edit UI in Customers tab (checkboxes)",
    ]),
    ("3.5  JARVIS Profile & Briefing", C_MID, [
        "Customer usage profile in chat.profiles",
        "Top 5 customers injected into every system prompt",
        "Auto-briefing card on Chat tab open (once per session)",
        "get_briefing agent tool — narrated morning summary",
        "Proactive P1 alert timer (15-min): ui.notify toast",
    ]),
]

cw = 4.22; ch = 2.72; sx = 0.28; gap_x = 0.075
for i, (name, col, bullets) in enumerate(milestones):
    row, c = divmod(i, 3)
    card(slide, sx + c*(cw+gap_x), HBAR+0.18 + row*ch,
         cw, ch-0.10, name, bullets, col, tsz=13, bsz=11)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 13 — Roadmap Phase 4
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Roadmap  —  Phase 4: AI-driven Interface",
       "Agent becomes the primary interface  ·  Data freshness is automatic  ·  docs/phase4-ai-driven-interface.md")

milestones = [
    ("4.1  MCP Tool Server", C_RED, [
        "All 29+ tools via FastMCP stdio transport",
        "Claude Desktop: ~/.claude/claude_desktop_config.json",
        "CB credentials via env vars — no hardcoding",
        "Long-running tool progress streaming",
        "Preflight tab: MCP server running / stopped status",
    ]),
    ("4.2  Agent-driven Freshness", C_BLUE, [
        "fetch_fresh_data(org) runs pipeline from agent turn",
        "Freshness thresholds: critical=4h · high=24h · normal=72h",
        "Prompt rule: 'call check_freshness before answering'",
        "Stale data ⚠ inline banner if refresh skipped",
        "Scrape progress shown in agent status strip",
    ]),
    ("4.3  PDF Report Generation", C_ORANGE, [
        "WeasyPrint (HTML→PDF, pure Python) preferred backend",
        "render_pdf_report(sections): markdown + chart + table",
        "Couchbase-branded template: logo · date · page numbers",
        "ECharts → PNG via pyecharts server-side render",
        "Chainlit: PDF sent as cl.File attachment",
    ]),
    ("4.4  Session Management", C_GREEN, [
        "Session picker: last 10 sessions from CB (sidebar)",
        "resume_session() — loads history + prior_session_block",
        "Topic tag extraction via lightweight LLM call",
        "Prior context chip at top of chat (collapsible)",
        "Chainlit thread sidebar: 'Resume in NiceGUI' deep-link",
    ]),
    ("4.5  Scheduled Pipeline", C_MID, [
        "pipeline_runner.py — headless, cron-compatible",
        "CLI: --org · --scope · --embed · --incremental",
        "Exit codes: 0=success · 1=partial · 2=total failure",
        "CB status doc (pipeline_status::current) for live UI",
        "Launchd / cron examples in docs/scheduling.md",
    ]),
    ("4.6 + 4.7  Chat-first & SFDC", C_RED, [
        "Chainlit: charts as cl.Image · scrape status inline card",
        "Salesforce: fetch_sfdc_accounts(se_email) SOQL query",
        "accounts CB collection: ARR · renewal · CSM · org_aliases",
        "org_aliases bridge SF account name → ticket org string",
        "renewal_risk: open P1/P2 + renewal within 90 days",
    ]),
]

cw = 4.22; ch = 2.72; sx = 0.28; gap_x = 0.075
for i, (name, col, bullets) in enumerate(milestones):
    row, c = divmod(i, 3)
    card(slide, sx + c*(cw+gap_x), HBAR+0.18 + row*ch,
         cw, ch-0.10, name, bullets, col, tsz=13, bsz=11)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 14 — Summary
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
header(slide, "Summary",
       "Supportal AI Platform v2.6.1  ·  Two apps, one shared library, one database")

# KPI row
kpis = [
    ("3",    "Applications",       C_RED),
    ("10",   "Shared Modules",     C_BLUE),
    ("29+",  "Agent Tools",        C_ORANGE),
    ("16.7K","Lines  (Strabo)",    C_MID),
    ("28",   "Pre-built Prompts",  C_GREEN),
    ("8+2",  "CB Indexes",         C_MUTEBL),
]
kw = 2.06; kgap = 0.08
for i, (val, lbl, col) in enumerate(kpis):
    kx = 0.28 + i*(kw+kgap)
    R(slide, kx, HBAR+0.18, kw, 1.10, col)
    T(slide, val, kx+0.10, HBAR+0.20, kw-0.20, 0.58,
      sz=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    T(slide, lbl, kx+0.10, HBAR+0.78, kw-0.20, 0.40,
      sz=11, color=C_WHITE, align=PP_ALIGN.CENTER)

card(slide, 0.28, HBAR+1.46, 6.18, 2.38, "What's Shipped  (v2.6.1)", [
    "✓ Cursus — single command, starts both UIs, watchdog + relay",
    "✓ Strabo — 7-tab NiceGUI dashboard, analytics, assets, Corax iframe",
    "✓ Corax — Chainlit chat, thread persistence, Top-K slider",
    "✓ Hybrid retrieval — FTS + Vector + N1QL → RRF + elbow Top-K",
    "✓ 29+ agent tools (query / intel / scrape / utility / time)",
    "✓ Docker: couchbase-init.sh idempotent bootstrap",
], C_RED, tsz=14, bsz=12)

card(slide, 6.74, HBAR+1.46, 6.31, 2.38, "What's Next  (Phase 3 → 4)", [
    "→ Fleet Dashboard — version distribution, risk heatmap",
    "→ Leading Indicators — at-risk clusters before tickets open",
    "→ Portfolio Management — named org groups with health rollups",
    "→ JARVIS Profile — ambient briefing, proactive P1 alerts",
    "→ MCP Tool Server — Claude Desktop as lightweight interface",
    "→ Salesforce Integration — blended support + ARR + renewal",
], C_BLUE, tsz=14, bsz=12)

R(slide, 0.28, HBAR+4.00, 12.77, 0.05, C_RED)

card(slide, 0.28, HBAR+4.10, 5.90, 2.00, "Documentation", [
    "README.md  ·  docs/architecture.html",
    "docs/workflow.html",
    "docs/phase3-fleet-analysis.md",
    "docs/phase4-ai-driven-interface.md",
    "docs/milestone-3.5-jarvis-profile.md",
], C_MID, tsz=13, bsz=11)

card(slide, 6.46, HBAR+4.10, 3.00, 2.00, "Tech Stack", [
    "Python 3.14  ·  NiceGUI 3.x",
    "Chainlit  ·  Couchbase SDK 4.x",
    "Playwright  ·  Claude (Anthropic)",
    "OpenAI  ·  Ollama  ·  LMStudio",
    "Docker  ·  Open Sans font",
], C_MID, tsz=13, bsz=11)

card(slide, 9.74, HBAR+4.10, 3.31, 2.00, "Naming Summary", [
    "Cursus  — Rome's imperial relay network",
    "Strabo  — Greek geographer of the ancient world",
    "Corax   — Greek: raven; founder of Western rhetoric",
    "supportal  — shared library, the connective tissue",
], C_ORANGE, tsz=13, bsz=11)


# ─────────────────────────────────────────────────────────────────────────────
out = HERE / "docs" / "Supportal_Architecture_v2.6.1.pptx"
prs.save(str(out))
print(f"✓  {out}")
print(f"   {len(prs.slides)} slides · Open Sans · Couchbase brand colors")
